"""PPTX export generator for dossier data.

Creates an investment research presentation with dark-themed slides.
All python-pptx imports are deferred to avoid startup crashes if not installed.
"""

from io import BytesIO
from typing import Any


def _safe_str(v: Any) -> str:
    if v is None:
        return "-"
    if isinstance(v, float):
        if abs(v) >= 1e12:
            return f"${v / 1e12:.2f}T"
        if abs(v) >= 1e9:
            return f"${v / 1e9:.2f}B"
        if abs(v) >= 1e6:
            return f"${v / 1e6:.1f}M"
        return f"{v:,.2f}"
    if isinstance(v, list):
        return ", ".join(str(x) for x in v[:5])
    if isinstance(v, dict):
        if "value" in v:
            return _safe_str(v["value"])
        if "score" in v:
            return _safe_str(v["score"])
        entries = [f"{k}: {_safe_str(val)}" for k, val in list(v.items())[:4]]
        return "; ".join(entries)
    return str(v)[:200]


def _flatten_for_slide(d: dict, max_items: int = 14) -> list[tuple[str, str]]:
    """Flatten dict to (label, value_str) pairs for slide display."""
    rows: list[tuple[str, str]] = []
    for k, v in d.items():
        if k.startswith("_"):
            continue
        label = k.replace("_", " ").title()
        if isinstance(v, dict):
            if "value" in v or "current" in v or "score" in v:
                rows.append((label, _safe_str(v)))
            else:
                for nk, nv in v.items():
                    if nk.startswith("_"):
                        continue
                    nested_label = f"{label} > {nk.replace('_', ' ').title()}"
                    rows.append((nested_label, _safe_str(nv)))
        else:
            rows.append((label, _safe_str(v)))
        if len(rows) >= max_items:
            break
    return rows[:max_items]


def generate_dossier_pptx(dossier_data: dict) -> bytes:
    """Generate a presentation from dossier data.

    Returns the PPTX as bytes suitable for StreamingResponse.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    # ── Design tokens ──
    BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
    TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT_SECONDARY = RGBColor(0x9C, 0xA3, 0xAF)
    TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
    ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
    ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
    ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
    ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)

    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def set_slide_bg(slide, color=BG_DARK):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title(slide, text, left=0.5, top=0.3, width=12, height=0.8, size=28):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT_WHITE
        p.font.bold = True
        return txBox

    def add_subtitle(slide, text, left=0.5, top=1.1, width=12, height=0.5, size=14):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT_SECONDARY
        return txBox

    def add_body_text(slide, text, left=0.5, top=1.8, width=12, height=5.0, size=12):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT_SECONDARY
        p.line_spacing = Pt(18)
        return txBox

    def add_bullet_list(slide, items, left=0.5, top=1.8, width=12, height=5.0, size=11, color=TEXT_SECONDARY):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items[:12]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"  {item}"
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.line_spacing = Pt(18)
            p.space_after = Pt(4)
        return txBox

    def add_accent_line(slide, left=0.5, top=1.55, width=2.0, color=ACCENT_BLUE):
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Pt(3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def add_kv_table(slide, data, left=0.5, top=1.8, width=12.0, row_height=0.35):
        rows = min(len(data), 14)
        if rows == 0:
            return
        table_shape = slide.shapes.add_table(rows, 2, Inches(left), Inches(top), Inches(width), Inches(rows * row_height))
        table = table_shape.table
        table.columns[0].width = Inches(4)
        table.columns[1].width = Inches(8)
        for i, (label, value) in enumerate(data[:rows]):
            cell_l = table.cell(i, 0)
            cell_l.text = str(label)
            for p in cell_l.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = TEXT_MUTED
                p.font.bold = True
            cell_v = table.cell(i, 1)
            cell_v.text = str(value) if value is not None else ""
            for p in cell_v.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = TEXT_WHITE
            for j in range(2):
                cell = table.cell(i, j)
                cell.fill.solid()
                if i % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(0x13, 0x16, 0x1B)
                else:
                    cell.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # ── Build presentation ──
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    content = dossier_data.get("content", {})
    entity = dossier_data.get("entity", "Unknown")
    confidence = dossier_data.get("confidence_score", 0)
    blank_layout = prs.slide_layouts[6]

    # ═══ SLIDE 1: Title ═══
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    add_subtitle(slide, "NANO BANA PRO  |  EQUITY RESEARCH", top=1.0, size=10)

    profile_name = ""
    profile = content.get("fundamentals", {}).get("profile", {})
    if isinstance(profile, dict):
        profile_name = profile.get("company_name", "")
    add_title(slide, profile_name or entity, top=2.0, size=40)
    add_subtitle(slide, entity, top=3.2, size=18)

    es = content.get("executive_summary", {})
    thesis = content.get("investment_thesis", {})
    rec = ""
    if isinstance(es, dict):
        rec = es.get("recommendation", "")
    if not rec and isinstance(thesis, dict):
        rec = thesis.get("recommendation", "")
    if rec:
        rec_box = add_subtitle(slide, rec.upper(), top=4.0, size=16)
        for p in rec_box.text_frame.paragraphs:
            r = rec.lower()
            if "buy" in r:
                p.font.color.rgb = ACCENT_GREEN
            elif "sell" in r:
                p.font.color.rgb = ACCENT_RED
            else:
                p.font.color.rgb = ACCENT_AMBER
            p.font.bold = True

    stats_parts = [f"Confidence: {confidence * 100:.0f}%"]
    if isinstance(thesis, dict):
        tp = thesis.get("target_price")
        cp = thesis.get("current_price")
        if tp:
            stats_parts.append(f"Target: ${tp:.2f}")
        if cp:
            stats_parts.append(f"Current: ${cp:.2f}")
    add_subtitle(slide, "  |  ".join(stats_parts), top=5.0, size=12)
    add_subtitle(slide, "CONFIDENTIAL  |  Generated by NanoBana AI", top=6.5, size=9)

    # ═══ SLIDE 2: Executive Summary ═══
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title(slide, "Executive Summary")
    add_accent_line(slide)

    if isinstance(es, dict):
        summary_text = es.get("executive_summary", "") or es.get("headline", "")
        if summary_text:
            add_body_text(slide, summary_text[:800], top=1.8, height=2.5, size=11)
        swot_items = []
        for category, prefix in [("strengths", "+"), ("weaknesses", "-"), ("opportunities", "+"), ("threats", "-")]:
            items = es.get(category, [])
            if isinstance(items, list):
                for item in items[:2]:
                    swot_items.append(f"[{category.upper()[:3]}] {prefix} {item}")
        if swot_items:
            add_bullet_list(slide, swot_items, top=4.5, height=2.5, size=10)

    # ═══ SLIDE 3: Key Financials ═══
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title(slide, "Key Financials")
    add_accent_line(slide)

    fundamentals = content.get("fundamentals", {})
    fin_data: list[tuple[str, str]] = []
    for section_key in ["profile", "market_data", "financials", "profitability"]:
        section = fundamentals.get(section_key, {})
        if isinstance(section, dict):
            fin_data.extend(_flatten_for_slide(section, max_items=4))
    if fin_data:
        add_kv_table(slide, fin_data[:14], top=1.8)

    # ═══ SLIDE 4: Valuation Summary ═══
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title(slide, "Valuation Analysis")
    add_accent_line(slide)

    valuation = content.get("valuation", {})
    val_data: list[tuple[str, str]] = []
    for method_key, method_label in [("dcf", "DCF"), ("comps", "Comps"), ("precedent_transactions", "Precedent Txns"), ("sum_of_parts", "SOTP")]:
        method = valuation.get(method_key, {})
        if isinstance(method, dict) and method:
            rows = _flatten_for_slide(method, max_items=3)
            for label, value in rows:
                val_data.append((f"[{method_label}] {label}", value))
    if val_data:
        add_kv_table(slide, val_data[:14], top=1.8)

    if isinstance(thesis, dict):
        scenarios = []
        for case_key, case_label in [("bull_case", "Bull"), ("base_case", "Base"), ("bear_case", "Bear")]:
            case = thesis.get(case_key, {})
            if isinstance(case, dict) and case:
                tp_val = case.get("target_price")
                prob = case.get("probability_pct")
                if tp_val:
                    scenarios.append(f"{case_label}: ${tp_val:.2f} ({prob}% prob)")
        if scenarios:
            add_bullet_list(slide, scenarios, top=5.5, height=1.5, size=10, color=ACCENT_BLUE)

    # ═══ SLIDE 5: Risk Overview ═══
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title(slide, "Risk Assessment")
    add_accent_line(slide, color=ACCENT_AMBER)

    risk_quality = content.get("risk_and_quality", {})
    risk = risk_quality.get("risk_assessment", {})
    risk_bullets: list[str] = []
    if isinstance(risk, dict):
        score = risk.get("overall_risk_score")
        tier = risk.get("risk_tier")
        if score is not None:
            risk_bullets.append(f"Overall Risk Score: {score}/100")
        if tier:
            risk_bullets.append(f"Risk Tier: {tier}")
        dims = risk.get("risk_dimensions", {})
        if isinstance(dims, dict):
            for dim_name, dim_data in dims.items():
                if isinstance(dim_data, dict):
                    risk_bullets.append(f"{dim_name.replace('_', ' ').title()}: {dim_data.get('score', 'N/A')}/100")
        key_risks = risk.get("key_risks", [])
        if isinstance(key_risks, list):
            for kr in key_risks[:4]:
                if isinstance(kr, dict):
                    risk_bullets.append(f"[{kr.get('severity', '?')}] {kr.get('risk', '?')}")
        red_flags = risk.get("red_flags", [])
        if isinstance(red_flags, list):
            for rf in red_flags[:3]:
                risk_bullets.append(f"RED FLAG: {rf}")
    if risk_bullets:
        add_bullet_list(slide, risk_bullets, top=1.8, size=11, color=ACCENT_AMBER)

    # ═══ SLIDE 6: Recommendation ═══
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title(slide, "Recommendation")
    add_accent_line(slide, color=ACCENT_GREEN)

    rec_items: list[str] = []
    if isinstance(thesis, dict):
        if thesis.get("recommendation"):
            rec_items.append(f"Recommendation: {thesis['recommendation']}")
        if thesis.get("conviction_level"):
            rec_items.append(f"Conviction: {thesis['conviction_level']}")
        if thesis.get("target_price"):
            rec_items.append(f"Target Price: ${thesis['target_price']:.2f}")
        if thesis.get("upside_pct"):
            pct = thesis["upside_pct"]
            rec_items.append(f"Upside: {'+' if pct > 0 else ''}{pct:.1f}%")
        if thesis.get("investment_thesis"):
            rec_items.append(f"Thesis: {thesis['investment_thesis'][:300]}")
        catalysts = thesis.get("catalysts", [])
        if isinstance(catalysts, list):
            for cat in catalysts[:4]:
                if isinstance(cat, dict):
                    rec_items.append(f"Catalyst: {cat.get('event', '?')} ({cat.get('timeline', '?')})")
    if rec_items:
        add_bullet_list(slide, rec_items, top=1.8, size=11, color=ACCENT_GREEN)

    add_subtitle(slide, "Generated by NanoBana AI. For informational purposes only. Not investment advice.", top=6.5, size=8)

    # ── Serialize ──
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
